import os
import random
import re

import google.generativeai as genai
from dotenv import load_dotenv
from pptx import Presentation

from Cache.default_prompt import prompt
from content_extractor import extract_contents_from_text

from layout_report_tool import supporting_parameters

# Load environment variables from .env file
load_dotenv()

# Get the API key from the environment
api_key = os.getenv("API_KEY")

# Configure the genai library with the API key
genai.configure(api_key=api_key)


def get_bot_response(topic, theme):
    user_text = topic
    text = ""
    input_string = user_text
    input_string = re.sub(r'[^\w\s.\-\(\)]', '', input_string)
    input_string = input_string.replace("\n", "")
    number = int(theme[-1])
    pptlink = None

    if number > 9:
        number = 1
        print("Unavailable design, using default design...")
    elif number < 1:
        number = 1
        print("Unavailable design, using default design...")
    else:
        print(f"Available design, using {theme} design...")

    # Generate a filename using OpenAI API
    filename_prompt = f"""Generate a short, descriptive filename based on the following input: \"{user_text}\".
Answer just with the short filename; no other explanation. 
Do not give extensions to files like my_file.txt. I just need a file name."""

    model = genai.GenerativeModel('gemini-pro')

    generation_config = genai.GenerationConfig(
        stop_sequences=None,
        temperature=0.9,
        top_p=1.0,
        top_k=32,
        candidate_count=1,
        max_output_tokens=400,
    )

    question = filename_prompt

    filename_response = model.generate_content(
        contents=question,
        generation_config=generation_config,
        stream=False,
    )

    filename = filename_response.text.strip().replace(" ", "_")

    if number <= 7:
        question = prompt(user_text)

        text = model.generate_content(
            contents=question,
            generation_config=generation_config,
            stream=False,
        )

        cache_dir = 'Powerpointer-main/Cache'
        if not os.path.exists(cache_dir):
            os.makedirs(cache_dir)

        with open(f'{cache_dir}/{filename}.txt', 'w', encoding='utf-8') as f:
            f.write(text.text)
        placeholder_indices_by_layout_, layout_indices_, _ = supporting_parameters(number)
        pptlink = create_ppt_default(f'{cache_dir}/{filename}.txt', number, filename)

    else:
        _, _, index_containing_placeholders = supporting_parameters(number)
        pptlink = create_ppt_custom(f'Powerpointer-main/Cache/custom_prompt.txt', number, filename,
                                    index_containing_placeholders)

    return pptlink, f'{cache_dir}/{filename}'


def create_ppt_custom(text_file, design_number, ppt_name, index_containing_placeholders):
    prs = Presentation(f"Powerpointer-main/Designs/Design-{design_number}.pptx")
    placeholder_index = index_containing_placeholders
    slide_layout_index = 0
    slide_layout_index_ = 0
    slide_count = 0
    header = ""
    content = ""
    last_slide_layout_index = -1
    firsttime = True

    with open(text_file, 'r', encoding='utf-8') as f:
        for line_num, line in enumerate(f):

            if line.startswith('#Title:'):
                header = line.replace('#Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = header
                continue

            elif line.startswith('#Slide:'):

                # Split the input string into lines and process
                with open(text_file, 'r', encoding='utf-8') as f:
                    lines = f.read()

                contents = extract_contents_from_text(lines)
                placeholder_content = []

                for idx, slide in enumerate(contents, 1):
                    placeholder_content.append(slide)

                content = placeholder_content

                slide = prs.slides.add_slide(prs.slide_layouts[slide_count + 1])
                title = slide.shapes.title
                title.text = header
                count = 0
                print(placeholder_index[slide_count + 1])
                print(len(content[slide_count]))

                for i in placeholder_index[slide_count + 1]:
                    body_shape = slide.shapes.placeholders[i]
                    tf = body_shape.text_frame
                    tf.text = content[slide_count][count]
                    count += 1

                content = ""
                slide_count += 1

            elif line.startswith('#Header:'):
                header = line.replace('#Header:', '').strip()
                continue

    prs.save(f'Powerpointer-main/GeneratedPresentations/{ppt_name}.pptx')

    ppt_path = str(f'Powerpointer-main/GeneratedPresentations/{ppt_name}.pptx')
    return ppt_path


def create_ppt_default(text_file, design_number, ppt_name):
    prs = Presentation(f"Powerpointer-main/Designs/Design-{design_number}.pptx")
    slide_count = 0
    header = ""
    content = ""
    last_slide_layout_index = -1
    firsttime = True

    layout_indices = [1, 7, 8]  # Valid slide layouts
    placeholder_indices = {1: 1, 7: 1, 8: 2}  # Corresponding placeholder indices

    with open(text_file, 'r', encoding='utf-8') as f:
        for line in f:
            if line.startswith('#Title:'):
                header = line.replace('#Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = header
                body_shape = slide.shapes.placeholders[1]
                continue
            elif line.startswith('#Slide:'):
                if slide_count > 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
                    title = slide.shapes.title
                    title.text = header
                    body_shape = slide.shapes.placeholders[slide_placeholder_index]
                    tf = body_shape.text_frame
                    tf.text = content
                content = ""
                slide_count += 1
                slide_layout_index = last_slide_layout_index

                while slide_layout_index == last_slide_layout_index:
                    if firsttime:
                        slide_layout_index = 1
                        slide_placeholder_index = 1
                        firsttime = False
                        break
                    slide_layout_index = random.choice(layout_indices)  # Select random slide layout index
                    slide_placeholder_index = placeholder_indices[slide_layout_index]

                last_slide_layout_index = slide_layout_index
                continue

            elif line.startswith('#Header:'):
                header = line.replace('#Header:', '').strip()
                continue

            elif line.startswith('#Content:'):
                content = line.replace('#Content:', '').strip()
                next_line = f.readline().strip()
                while next_line and not next_line.startswith('#'):
                    content += '\n' + next_line
                    next_line = f.readline().strip()
                continue

    # Add the last slide if there is content remaining
    if content:
        slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
        title = slide.shapes.title
        title.text = header
        body_shape = slide.shapes.placeholders[slide_placeholder_index]
        tf = body_shape.text_frame
        tf.text = content

    prs.save(f'Powerpointer-main/GeneratedPresentations/{ppt_name}.pptx')
    ppt_path = str(f'Powerpointer-main/GeneratedPresentations/{ppt_name}.pptx')
    return ppt_path
