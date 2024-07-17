from pptx import Presentation
import pandas as pd


def list_placeholders(design_number):
    placeholders_data = []

    prs = Presentation(f"Powerpointer-main/Designs/Design-{design_number}.pptx")
    for i, layout in enumerate(prs.slide_layouts):
        for shape in layout.placeholders:
            placeholder_format = shape.placeholder_format
            placeholders_data.append([
                i,
                shape.placeholder_format.idx,
                shape.shape_id,
                shape.name,
                placeholder_format.type,
                shape.width,
                shape.height,
                shape.left,
                shape.top
            ])

    # Create DataFrame
    df = pd.DataFrame(placeholders_data, columns=[
        'Layout Index', 'Placeholder Index', 'Shape ID', 'Name', 'Type', 'Width', 'Height', 'Left', 'Top'
    ])

    return df


def color_rows_by_layout(df):
    def apply_colors(row):
        colors = ['background-color: #ffcccc',  # Light red
                  'background-color: #ccffcc',  # Light green
                  'background-color: #ccccff',  # Light blue
                  'background-color: #ffffcc',  # Light yellow
                  'background-color: #ffccff',  # Light pink
                  'background-color: #ccffff',  # Light cyan
                  'background-color: #ffd700',  # Gold
                  'background-color: #e6e6fa']  # Lavender
        return [colors[row['Layout Index'] % len(colors)]] * len(row)

    return df.style.apply(apply_colors, axis=1)


def get_placeholder_indices_by_layout(df):
    placeholder_indices_by_layout = df.groupby('Layout Index')['Placeholder Index'].apply(list).tolist()
    layout_indices = df['Layout Index'].unique().tolist()

    # Filter and group the DataFrame
    filtered_df = df[df['Name'].str.startswith('Content Placeholder')]
    grouped = filtered_df.groupby('Layout Index')['Placeholder Index'].apply(list)

    # Convert to list of lists
    index_containing_placeholders = grouped.tolist()

    return placeholder_indices_by_layout, layout_indices, index_containing_placeholders


def supporting_parameters(design_numer):
    design_number = design_numer
    placeholders_df = list_placeholders(design_number)

    # Display the styled DataFrame
    styled_df = color_rows_by_layout(placeholders_df)

    # Get placeholder indices by layout
    placeholder_indices_by_layout_, layout_indices_, index_containing_placeholders = get_placeholder_indices_by_layout(
        placeholders_df)

    return placeholder_indices_by_layout_, layout_indices_, index_containing_placeholders
